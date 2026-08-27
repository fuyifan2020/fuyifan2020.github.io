#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
convert_source.py —— 把各种格式的文章源，提取正文并生成带 front-matter 的草稿 .md

用法:
    python3 .codebuddy/skills/post/scripts/convert_source.py <源文件路径> [--out Blogs/2026_8_27]

说明:
    - 仅依赖 Python 标准库即可处理 .txt / .md / .html。
    - .pdf / .docx / .pptx 若装了对应工具或库会自动使用；缺失时打印提示并退而求其次，
      不会中断（交给人工补正）。
    - 生成的草稿含 front-matter 模板，需人工精修标题/分类/摘要与正文结构后再跑发布脚本。
"""
import os
import re
import sys
import shutil
import argparse
import subprocess
import pathlib
from datetime import date

ROOT = str(pathlib.Path(__file__).resolve().parents[4])


def slugify(s):
    s = s.strip().lower()
    s = re.sub(r"\s+", "-", s)
    s = re.sub(r"[^a-z0-9\-一-鿿]", "", s)
    return s.strip("-") or "post"


def extract_html(path):
    from html.parser import HTMLParser

    class P(HTMLParser):
        def __init__(self):
            super().__init__()
            self.out = []
            self.skip = 0

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style", "head"):
                self.skip += 1
            if tag in ("p", "div", "br", "li", "h1", "h2", "h3", "h4", "tr"):
                self.out.append("\n")

        def handle_endtag(self, tag):
            if tag in ("script", "style", "head") and self.skip:
                self.skip -= 1

        def handle_data(self, data):
            if self.skip == 0:
                t = data.strip()
                if t:
                    self.out.append(t + " ")

    p = P()
    p.feed(open(path, encoding="utf-8", errors="ignore").read())
    return re.sub(r"\n{2,}", "\n\n", "".join(p.out)).strip()


def extract_pdf(path):
    if shutil.which("pdftotext"):
        try:
            return subprocess.check_output(["pdftotext", path, "-"], text=True)
        except Exception:
            pass
    print("[convert] 未找到 pdftotext（poppler）。可 `conda install poppler` / "
          "系统装 poppler-utils，或用 pandoc；否则请手动粘贴 PDF 文字。", file=sys.stderr)
    return ""


def extract_docx(path):
    try:
        import docx
        d = docx.Document(path)
        return "\n\n".join(p.text for p in d.paragraphs if p.text.strip())
    except ImportError:
        print("[convert] 未安装 python-docx；可 `pip install python-docx` 或 "
              "`pandoc in.docx -t markdown`。", file=sys.stderr)
        return ""


def extract_pptx(path):
    try:
        from pptx import Presentation
        out = []
        for slide in Presentation(path).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        out.append(t)
        return "\n\n".join(out)
    except ImportError:
        print("[convert] 未安装 python-pptx；可 `pip install python-pptx` 或 "
              "`pandoc in.pptx -t markdown`。", file=sys.stderr)
        return ""


def extract(path):
    ext = path.lower().rsplit(".", 1)[-1]
    if ext in ("txt",):
        return open(path, encoding="utf-8", errors="ignore").read().strip()
    if ext in ("html", "htm"):
        return extract_html(path)
    if ext == "pdf":
        return extract_pdf(path).strip()
    if ext == "docx":
        return extract_docx(path).strip()
    if ext == "pptx":
        return extract_pptx(path).strip()
    # 未知扩展名：尽力当纯文本
    try:
        return open(path, encoding="utf-8", errors="ignore").read().strip()
    except Exception:
        return ""


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("source", help="源文件路径")
    ap.add_argument("--out", default=None, help="输出 Blogs/<日期> 目录，默认按今天日期")
    a = ap.parse_args()

    src = a.source
    if not os.path.isfile(src):
        print("源文件不存在:", src)
        sys.exit(1)

    ext = src.lower().rsplit(".", 1)[-1]
    base = os.path.splitext(os.path.basename(src))[0]
    slug = slugify(base)
    today = date.today().strftime("%Y-%m-%d")
    outdir = a.out or os.path.join("Blogs", date.today().strftime("%Y_%m_%d"))
    os.makedirs(outdir, exist_ok=True)
    out_md = os.path.join(outdir, slug + ".md")

    # 已是带 front-matter 的 .md：直接拷贝，避免重复包裹
    raw = open(src, encoding="utf-8", errors="ignore").read()
    if ext == "md" and raw.lstrip().startswith("---"):
        shutil.copy(src, out_md)
        print("[完成] 检测到现成 front-matter，已拷贝:", out_md)
        print("   请人工精修后再运行 publish_post.py")
        return

    text = extract(src)
    if not text:
        print("[warn] 未能提取到正文，已生成占位草稿，请手动整理。", file=sys.stderr)

    fm = (
        "---\n"
        f"title: {base}\n"
        f"slug: {slug}\n"
        f"date: {today}\n"
        "category: 技术\n"
        "excerpt: \n"
        "paper: \n"
        "---\n\n"
    )
    body = text if text else "(在此粘贴 / 整理正文)"
    with open(out_md, "w", encoding="utf-8") as f:
        f.write(fm + body + "\n")

    print("[完成] 草稿已生成:", out_md)
    print("   请人工精修标题 / 分类 / 摘要与正文结构，再运行 publish_post.py")


if __name__ == "__main__":
    main()
